#!/usr/bin/env python3
"""mineru2pptx.py - MinerU 输出 -> 可编辑 PPT (纯图片+文字,无背景)

用法:
  mineru -p <input_dir> -o <mineru_output_dir> --backend pipeline
  python3 mineru2pptx.py <input_dir> <mineru_output_dir> <output.pptx> [--slide-width 13.333]

字号: 三档分档 (大 30pt / 中 26pt / 小 22pt)
行距: 1.3 倍
文本框: 自动扩展宽度 (避免挤压文字到下一行)
"""
import argparse
import json
import logging
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Emu, Pt

logger = logging.getLogger(__name__)


# 文本框扩展参数: 让文字框略宽, 避免挤压但不过度
TEXTBOX_PADDING_PX = 8       # 文本框四周扩展像素 (轻微 padding)
TEXTBOX_MIN_WIDTH_PX = 80    # 文本框最小宽度(像素)


def build_clean_pptx(jpeg_dir, mineru_dir, output_pptx, slide_w_inches=13.333, snap_to_pool=False):
    from font_estimator import LINE_SPACING, estimate_font_size

    prs = Presentation()
    blank = prs.slide_layouts[6]

    jpeg_dir = Path(jpeg_dir)
    mineru_dir = Path(mineru_dir)

    aspect = 16 / 9
    sw = int(slide_w_inches * 914400)
    sh = int(sw / aspect)

    image_ids = sorted([
        int(p.stem)
        for p in jpeg_dir.iterdir()
        if p.suffix.lower() in (".jpg", ".jpeg") and p.stem.isdigit()
    ])

    for img_id in image_ids:
        cl_path = mineru_dir / str(img_id) / "auto" / f"{img_id}_content_list.json"
        if not cl_path.exists():
            logger.debug("skip %s: not found", img_id)
            continue

        img = Image.open(jpeg_dir / f"{img_id}.jpeg")
        W, H = img.size

        prs.slide_width = Emu(sw)
        prs.slide_height = Emu(sh)
        slide = prs.slides.add_slide(blank)

        with open(cl_path) as f:
            blocks = json.load(f)

        for blk in blocks:
            btype = blk.get("type")
            bbox_norm = blk.get("bbox", [])

            if btype in ("page_header", "page_footer") and not blk.get("text", "").strip():
                continue
            if not bbox_norm or len(bbox_norm) != 4:
                continue

            x1 = bbox_norm[0] / 1000 * W
            y1 = bbox_norm[1] / 1000 * H
            x2 = bbox_norm[2] / 1000 * W
            y2 = bbox_norm[3] / 1000 * H
            w_px = x2 - x1
            h_px = y2 - y1

            if btype == "image":
                emu_x = Emu(int(x1 / W * sw))
                emu_y = Emu(int(y1 / H * sh))
                emu_w = Emu(int(w_px / W * sw))
                emu_h = Emu(int(h_px / H * sh))

                img_filename = Path(blk.get("img_path", "")).name
                cropped_img = mineru_dir / str(img_id) / "auto" / "images" / img_filename
                if cropped_img.exists():
                    slide.shapes.add_picture(str(cropped_img), emu_x, emu_y, emu_w, emu_h)

            elif btype in ("text", "title", "paragraph", "list", "page_header", "page_footer"):
                text = blk.get("text", "").strip()
                if not text:
                    continue

                font_pt = estimate_font_size(text, w_px, h_px, snap=snap_to_pool)

                # 文本框: 轻微 padding (4 边各 +8px), 保证 min width
                cx = max(0, x1 - TEXTBOX_PADDING_PX)
                cy = max(0, y1 - TEXTBOX_PADDING_PX)
                cw = min(W - cx, w_px + TEXTBOX_PADDING_PX * 2)
                ch = min(H - cy, h_px + TEXTBOX_PADDING_PX * 2)
                cw = max(cw, TEXTBOX_MIN_WIDTH_PX)

                emu_x = Emu(int(cx / W * sw))
                emu_y = Emu(int(cy / H * sh))
                emu_w = Emu(int(cw / W * sw))
                emu_h = Emu(int(ch / H * sh))

                tb = slide.shapes.add_textbox(emu_x, emu_y, emu_w, emu_h)
                tf = tb.text_frame
                tf.word_wrap = True
                tf.margin_left = tf.margin_right = Emu(0)
                tf.margin_top = tf.margin_bottom = Emu(0)
                p = tf.paragraphs[0]
                p.text = text
                p.font.size = Pt(font_pt)
                p.line_spacing = LINE_SPACING
                if btype == "title":
                    p.font.bold = True

    prs.save(output_pptx)
    logger.info("saved: %s", output_pptx)
    logger.info("  slides: %d", len(prs.slides))


def _parse_args(argv=None):
    parser = argparse.ArgumentParser(
        description="MinerU 输出 (auto/*_content_list.json) → 可编辑 PPTX",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=__doc__,
    )
    parser.add_argument("input_dir", help="输入图片目录 (含 *.jpg 或 *.jpeg)")
    parser.add_argument("mineru_dir", help="MinerU 输出目录 (含 auto/*_content_list.json)")
    parser.add_argument("output_pptx", help="输出的 .pptx 文件")
    parser.add_argument(
        "--slide-width",
        type=float,
        default=13.333,
        help="slide 宽度(英寸), 默认 13.333 (16:9)",
    )
    parser.add_argument(
        "--log-level",
        default="INFO",
        choices=["DEBUG", "INFO", "WARNING", "ERROR"],
        help="日志级别 (默认 INFO)",
    )
    parser.add_argument(
        "--snap",
        action="store_true",
        help="把字号 snap 到候选池 (常见 PPT 字号离散值, 减少 ±1-2pt 误差)",
    )
    return parser.parse_args(argv)


if __name__ == "__main__":
    args = _parse_args()
    logging.basicConfig(
        level=getattr(logging, args.log_level),
        format="%(asctime)s [%(levelname)s] %(message)s",
    )
    if args.snap:
        logger.info("字号 snap to pool 已启用 (CANDIDATE_POOL)")
    build_clean_pptx(args.input_dir, args.mineru_dir, args.output_pptx, args.slide_width, snap_to_pool=args.snap)
